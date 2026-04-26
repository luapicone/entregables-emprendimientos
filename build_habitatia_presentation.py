from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = 'Presentacion_HabitatIA_5_7_minutos.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    'navy': RGBColor(15, 23, 42),
    'blue': RGBColor(37, 99, 235),
    'sky': RGBColor(14, 165, 233),
    'orange': RGBColor(249, 115, 22),
    'green': RGBColor(22, 163, 74),
    'red': RGBColor(220, 38, 38),
    'slate': RGBColor(71, 85, 105),
    'light': RGBColor(248, 250, 252),
    'muted': RGBColor(226, 232, 240),
    'dark_text': RGBColor(30, 41, 59),
}


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, dark=False):
    color = COLORS['light'] if dark else COLORS['navy']
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.8), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = color
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.72), Inches(1.15), Inches(11), Inches(0.5))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = 'Aptos'
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLORS['muted'] if dark else COLORS['slate']


def add_footer(slide, text, dark=False):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = 'Aptos'
    r.font.size = Pt(9)
    r.font.color.rgb = COLORS['muted'] if dark else COLORS['slate']


def add_bullets(slide, items, left=0.9, top=1.7, width=6.0, height=4.8, font_size=19, dark=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        p.bullet = True
        r = p.add_run()
        r.text = item
        r.font.name = 'Aptos'
        r.font.size = Pt(font_size)
        r.font.color.rgb = COLORS['light'] if dark else COLORS['dark_text']
    return box


def add_callout(slide, title, value, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.bold = True
    r1.font.size = Pt(24)
    r1.font.color.rgb = COLORS['light']
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    r2.font.size = Pt(11)
    r2.font.color.rgb = COLORS['light']


def add_process_boxes(slide, labels, left=0.7, top=2.0, total_width=12.0):
    gap = 0.18
    width = (total_width - gap * (len(labels)-1)) / len(labels)
    for i, label in enumerate(labels):
        x = left + i * (width + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(width), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['light'] if i % 2 == 0 else COLORS['muted']
        shape.line.color.rgb = COLORS['blue']
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = 'Aptos'
        r.font.bold = True
        r.font.size = Pt(14)
        r.font.color.rgb = COLORS['navy']
        if i < len(labels)-1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + width + 0.02), Inches(top + 0.48), Inches(0.12), Inches(0.42))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['orange']
            arrow.line.color.rgb = COLORS['orange']


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['navy'])
add_title(slide, 'HabitatIA', 'Presentación ejecutiva | 5 a 7 minutos | Emprendimientos', dark=True)

box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(7.7), Inches(3.3))
tf = box.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = 'Planificación habitacional asistida + reutilización inteligente de materiales'
r.font.name = 'Aptos Display'
r.font.bold = True
r.font.size = Pt(28)
r.font.color.rgb = COLORS['light']

p2 = tf.add_paragraph()
p2.space_before = Pt(18)
r2 = p2.add_run()
r2.text = 'Una plataforma PropTech que ordena la decisión inicial de construir, reduce incertidumbre y abre ahorro real mediante un marketplace de sobrantes de obra.'
r2.font.name = 'Aptos'
r2.font.size = Pt(19)
r2.font.color.rgb = COLORS['muted']

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.8), Inches(3.1), Inches(3.4))
shape.fill.solid(); shape.fill.fore_color.rgb = COLORS['orange']; shape.line.color.rgb = COLORS['orange']
tf = shape.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Problema\nreal'; r.font.bold = True; r.font.size = Pt(24); r.font.color.rgb = COLORS['light']
for txt in ['Falta de claridad', 'Sobrecostos', 'Materiales sobrantes mal usados']:
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; r.font.size = Pt(14); r.font.color.rgb = COLORS['light']
add_footer(slide, 'Equipo: Luca Picone, Ignacio Sanguinetti, Aquiles Luzuriaga, Antonio Cocca y Felipe Villares', dark=True)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
add_title(slide, '1. Problema y oportunidad')
add_bullets(slide, [
    'Muchas familias quieren construir o ampliar, pero se frenan antes de empezar por falta de claridad presupuestaria, técnica y visual.',
    'La etapa de preplanificación hoy está fragmentada: ideas sueltas, referencias dispersas, materiales poco entendibles y presupuestos inciertos.',
    'En paralelo, arquitectos, estudios y constructoras acumulan sobrantes con valor de uso que suelen inmovilizarse, venderse mal o desperdiciarse.',
    'La oportunidad está en integrar ambas fricciones en un solo flujo: ordenar la decisión y capturar ahorro real.'
], width=7.0, font_size=17)
add_callout(slide, 'Déficit habitacional estimado en Argentina', '3,24 M', 8.4, 1.8, 1.8, 1.4, COLORS['blue'])
add_callout(slide, 'Hogares urbanos con algún problema habitacional', '10,7 M', 10.4, 1.8, 1.8, 1.4, COLORS['sky'])
add_callout(slide, 'Hogares propietarios de vivienda y terreno (INDEC 2S 2023)', '60,2 %', 8.4, 3.5, 3.8, 1.4, COLORS['green'])
add_callout(slide, 'Desperdicio/ excedente conservador de materiales en obra', '10-15 %', 8.4, 5.2, 3.8, 1.4, COLORS['orange'])
add_footer(slide, 'Fuentes del entregable: Cámara Argentina de la Construcción, Fundación Tejido Urbano, INDEC.')

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(239, 246, 255))
add_title(slide, '2. Solución propuesta: qué hace HabitatIA')
add_process_boxes(slide, [
    '1. Usuario carga\nnecesidades, presupuesto\ny preferencias',
    '2. IA + reglas\nordenan la información\ny generan propuesta',
    '3. Render + plano\norientativo + cómputo\ny costo preliminar',
    '4. Marketplace\ncompara materiales\ny detecta ahorro'
])
add_bullets(slide, [
    'HabitatIA no vende solo una imagen atractiva: traduce una intención difusa en una propuesta accionable.',
    'El valor fuerte está en combinar visualización, lectura espacial, materiales, costo estimado y ahorro potencial.',
    'La arquitectura planteada es modular: LLMs para interpretar, reglas para coherencia, datos estructurados para materiales y matching para el marketplace.'
], top=4.1, width=11.3, font_size=17)
add_footer(slide, 'Promesa central: claridad inicial y mejor decisión, no reemplazo total del arquitecto ni presupuesto ejecutivo final.')

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
add_title(slide, '3. Qué validamos y qué aprendimos')
add_bullets(slide, [
    'El dolor principal en particulares y familias no es “no tener ideas”, sino no poder calcular costos, materiales y alternativas con seguridad.',
    'El render por sí solo no alcanza. La propuesta gana fuerza cuando se combina con plano orientativo, listado de materiales y presupuesto preliminar.',
    'El marketplace interesa, pero solo si resuelve confianza: reputación, evidencia visual, garantías y logística ordenada.',
    'Del lado profesional, la adopción depende de simplicidad operativa, baja carga administrativa y salida rápida del excedente.',
    'La narrativa comercial no debe apoyarse solo en IA. Tiene que apoyarse en utilidad, claridad y respaldo.'
], width=11.2, font_size=17)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(6.2), Inches(11.8), Inches(0.65))
bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['navy']; bar.line.color.rgb = COLORS['navy']
tf = bar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Insight clave: el núcleo del valor no es la IA, es reducir incertidumbre antes de gastar mal.'; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = COLORS['light']
add_footer(slide, 'Base: entrevistas y customer discovery incorporados en el entregable final.')

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(248, 250, 252))
add_title(slide, '4. Modelo de negocio y estrategia de entrada')
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.8), Inches(3.7))
left_box.fill.solid(); left_box.fill.fore_color.rgb = RGBColor(230, 244, 255); left_box.line.color.rgb = COLORS['blue']
tf = left_box.text_frame
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Monetización por fases'; r.font.bold = True; r.font.size = Pt(22); r.font.color.rgb = COLORS['navy']
for txt in [
    '1. Entrada gratuita para probar la experiencia inicial.',
    '2. Plan premium sugerido alrededor de ARS 29.900 para cómputo detallado y reporte ampliado.',
    '3. Comisión de 8 % sobre transacciones del marketplace cuando exista matching real.',
    '4. A futuro, servicios B2B, visibilidad y alianzas.'
]:
    p = tf.add_paragraph(); p.bullet = True; p.space_after = Pt(8)
    r = p.add_run(); r.text = txt; r.font.size = Pt(16); r.font.color.rgb = COLORS['dark_text']

right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.7), Inches(5.6), Inches(3.7))
right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(240, 253, 244); right_box.line.color.rgb = COLORS['green']
tf = right_box.text_frame
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Go to market inicial'; r.font.bold = True; r.font.size = Pt(22); r.font.color.rgb = COLORS['navy']
for txt in [
    'Foco en hogares con intención concreta de construir o ampliar y alta sensibilidad al costo.',
    'Segundo frente: personas con terreno propio.',
    'Tercer frente: arquitectos y estudios con sobrantes para poblar el marketplace.',
    'Canales sugeridos: sitio, landing pages, WhatsApp, contenidos y alianzas con corralones/profesionales.'
]:
    p = tf.add_paragraph(); p.bullet = True; p.space_after = Pt(8)
    r = p.add_run(); r.text = txt; r.font.size = Pt(16); r.font.color.rgb = COLORS['dark_text']
add_footer(slide, 'La lógica es cobrar cuando el usuario percibe profundidad y cuando la plataforma genera valor transaccional real.')

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
add_title(slide, '5. Riesgos críticos y condiciones de éxito')
add_bullets(slide, [
    'Riesgo 1: que el producto sea percibido como generador de imágenes sin profundidad económica o funcional.',
    'Riesgo 2: que el usuario interprete el presupuesto preliminar como cifra final y se rompa la confianza.',
    'Riesgo 3: que el marketplace no alcance masa crítica suficiente en oferta y calidad.',
    'Riesgo 4: dispersarse demasiado rápido e intentar resolver toda la cadena de construcción desde el inicio.'
], width=7.0, font_size=17)
add_callout(slide, 'Condición de éxito #1', 'Resolver bien la preplanificación', 8.2, 1.8, 3.9, 1.2, COLORS['blue'])
add_callout(slide, 'Condición de éxito #2', 'Construir confianza progresiva', 8.2, 3.3, 3.9, 1.2, COLORS['green'])
add_callout(slide, 'Condición de éxito #3', 'Activar marketplace con foco en calidad', 8.2, 4.8, 3.9, 1.2, COLORS['orange'])
add_footer(slide, 'La clave no es escalar rápido, sino demostrar utilidad, credibilidad y ahorro verificable.')

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['navy'])
add_title(slide, '6. Cierre: por qué HabitatIA tiene sentido', dark=True)
add_bullets(slide, [
    'Ataca un dolor real y frecuente: decidir cómo construir sin información clara al inicio.',
    'Integra dos capas con lógica económica: planificación asistida + reutilización inteligente de materiales.',
    'Tiene una oportunidad grande, pero una entrada realista: empezar por claridad y ahorro, no por escala total.',
    'Si valida confianza, precisión percibida y masa crítica mínima, puede convertirse en una PropTech con valor social, económico y ambiental.'
], width=10.8, font_size=19, dark=True)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(4.95), Inches(3.5), Inches(1.1))
shape.fill.solid(); shape.fill.fore_color.rgb = COLORS['orange']; shape.line.color.rgb = COLORS['orange']
tf = shape.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'HabitatIA organiza una decisión que hoy llega desordenada.'; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = COLORS['light']
add_footer(slide, 'Fin | Presentación sugerida para exposición oral breve (5 a 7 minutos).', dark=True)

prs.save(OUTPUT)
print(f'OK: {OUTPUT}')
